from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parent
OUTPUT_PATH = PROJECT_ROOT / "docs" / "phase4_presentation_professional.pptx"

TITLE_COLOR = RGBColor(12, 24, 38)
ACCENT_COLOR = RGBColor(198, 128, 44)
BODY_COLOR = RGBColor(28, 39, 52)
MUTED_COLOR = RGBColor(92, 105, 118)
WHITE = RGBColor(255, 255, 255)
SOFT_BG = RGBColor(246, 248, 251)


SLIDES = [
    {
        "title": "Developing and Optimizing Data Structures for Real-World Applications Using Python",
        "subtitle": "Search engine optimization with an inverted index, trie, and top-k heap",
        "bullets": [],
        "notes": "This project applies data-structure design to a search engine context. The final system shows how indexing, prefix suggestions, and ranking can be combined in a modular Python implementation.",
        "script": "Today I am presenting a search engine project built around three data structures: an inverted index, a trie, and a top-k heap. The goal was not only to implement the structures, but to optimize them so they remain efficient when documents are added, removed, and queried repeatedly.",
        "is_title": True,
    },
    {
        "title": "Problem and Context",
        "bullets": [
            "Search engines need fast retrieval over large text corpora.",
            "Updates must be efficient when documents are added or removed.",
            "Users expect ranked results and prefix suggestions.",
        ],
        "notes": "The project focuses on the core operations a search engine performs most often: storing documents, finding matching documents quickly, and ranking the most relevant results.",
        "script": "The application context is search because it is one of the clearest real-world examples of data structures driving performance. Search systems need to retrieve documents quickly, update their index when content changes, and return ranked results that feel responsive to users.",
    },
    {
        "title": "Data Structure Choices",
        "bullets": [
            "Inverted index for term-to-document lookup.",
            "Trie for prefix search and autocomplete.",
            "Top-k heap for efficient result selection.",
        ],
        "notes": "Each structure solves a different part of the problem. The inverted index handles retrieval, the trie supports vocabulary suggestions, and the heap keeps ranking efficient without sorting every candidate.",
        "script": "The inverted index maps terms to documents, which makes retrieval efficient. The trie supports prefix suggestions, which is useful for autocomplete-style interactions. The top-k heap keeps only the best results, so the ranking stage does not waste time sorting every candidate document.",
    },
    {
        "title": "Phase 1 Design Rationale",
        "bullets": [
            "Hash tables provide fast average-case lookups.",
            "Tries naturally support prefix traversal.",
            "Min-heaps are ideal for top-k selection.",
        ],
        "notes": "The design was chosen for fit, not novelty. The goal was to select structures that map cleanly to search-engine operations and are practical to implement in Python.",
        "script": "In phase 1, the focus was the design rationale. The key decision was to choose structures that match the workload. Hash tables are ideal for term lookup, tries are ideal for prefix traversal, and heaps are ideal for selecting the best few results.",
    },
    {
        "title": "Phase 2 Proof of Concept",
        "bullets": [
            "Added document insertion and removal.",
            "Supported AND and OR queries.",
            "Added basic prefix suggestions.",
        ],
        "notes": "The proof of concept proved the architecture worked end to end. It demonstrated that the application could index text, search it, and return ranked results in a modular way.",
        "script": "Phase 2 was the proof of concept. At that point, the code proved the workflow worked end to end: documents could be added, searched, and removed. This phase mattered because it confirmed the design before any optimization work was added.",
    },
    {
        "title": "Phase 3 Optimizations",
        "bullets": [
            "Incremental trie updates through DocumentDelta.",
            "Targeted deletion using per-document term maps.",
            "Smallest-posting-list-first intersection for AND queries.",
            "IDF caching with version-based invalidation.",
        ],
        "notes": "The main improvement in phase 3 was reducing repeated work. Instead of rebuilding structures after every mutation, the code updates only what changed and reuses cached relevance statistics when safe.",
        "script": "Phase 3 is where the design becomes more scalable. The trie is updated incrementally, deletion touches only affected postings, AND queries start with the smallest posting lists, and TF-IDF scores reuse cached IDF values. The common pattern is avoiding unnecessary work.",
    },
    {
        "title": "Ranking Methods",
        "bullets": [
            "TF-IDF is simple and interpretable.",
            "BM25 adds length normalization and saturation.",
            "Both are supported in the final engine.",
        ],
        "notes": "Supporting both ranking strategies makes the system more flexible and provides a useful comparison point in the evaluation section.",
        "script": "The project supports both TF-IDF and BM25. TF-IDF is easy to interpret and a good baseline. BM25 is more robust for ranking because it includes term saturation and document-length normalization. Keeping both in the system makes the evaluation more complete.",
    },
    {
        "title": "Testing and Validation",
        "bullets": [
            "Unit tests cover search, deletion, suggestions, and invalid input.",
            "Scaling tests use a larger synthetic corpus.",
            "BM25 and TF-IDF are both exercised.",
            "Results are verified through automated execution, not manual inspection.",
        ],
        "notes": "The test suite is designed to show correctness and robustness, not just happy-path behavior. The scaling test is especially important because it confirms the design still works when the input size grows.",
        "script": "The tests cover the main behaviors that matter: search, deletion, prefix lookup, ranking, invalid inputs, and large-dataset behavior. That gives confidence that the system is not just working for one small example but remains stable under growth and mutation.",
    },
    {
        "title": "Test Execution Results",
        "bullets": [
            "Command: python -m unittest discover -s tests -p \"test_*.py\" -v",
            "Outcome: 17 tests passed, 0 failed.",
            "Runtime: approximately 0.05 seconds.",
            "Coverage files: test_data_structures.py, test_search_engine.py, test_phase3_scaling.py.",
        ],
        "notes": "This slide reports the actual automated test run performed for phase 4. The goal is to show verifiable evidence that the implementation works across functionality, ranking, updates, and scaling behavior.",
        "script": "These are the actual test execution results from the project. I ran the full unittest discovery command across the tests folder. The suite completed with 17 tests passed and zero failures in about 0.05 seconds. The tests include core data structures, end-to-end search engine behavior, and scaling coverage.",
    },
    {
        "title": "Performance Discussion",
        "bullets": [
            "Incremental updates reduce rebuild cost.",
            "Targeted deletion avoids full scans.",
            "Caching improves repeated query performance.",
        ],
        "notes": "The major performance gains come from avoiding unnecessary work. That is the central theme of the optimization phase.",
        "script": "The performance story is simple. The optimized version reduces rebuild work, narrows deletion work, and caches repeated scoring values. Those changes are enough to improve responsiveness without making the implementation overly complicated.",
    },
    {
        "title": "Algorithm Overview and Sample Output",
        "bullets": [],
        "notes": "This slide summarizes the query pipeline and shows the kind of ranked output the system produces on the demo dataset.",
        "script": "The algorithm can be summarized in three steps. First, tokenize the query and collect candidate documents. Second, score only those candidates using the selected ranking function. Third, keep only the top k results in a heap. On the demo data, a query like search index returns doc-002, while search graph returns doc-001, doc-002, and doc-004.",
    },
    {
        "title": "Strengths and Limitations",
        "bullets": [
            "Strong modularity and clear code organization.",
            "Efficient enough for a classroom-scale search system.",
            "No compression, persistence, or distributed indexing.",
        ],
        "notes": "The final solution is credible and easy to explain, but it is still a simplified model of a production search engine.",
        "script": "The implementation is intentionally compact, which is a strength for a project submission. At the same time, it is not yet a production search engine. It does not have compression, persistence, or distributed indexing, so those are the main limitations to acknowledge.",
    },
    {
        "title": "Future Work",
        "bullets": [
            "Add phrase queries and metadata filters.",
            "Compress postings lists.",
            "Persist the index and add benchmark charts.",
        ],
        "notes": "These extensions would move the project closer to a production-quality retrieval system and create a stronger basis for further research.",
        "script": "The next steps are clear. A more advanced version could add phrase search, metadata filtering, compressed postings, and persistence. Those changes would turn this from a strong classroom prototype into a more realistic retrieval engine.",
    },
    {
        "title": "Closing",
        "bullets": [
            "Data structure choice drives search performance.",
            "Incremental maintenance improves scalability.",
            "The final design balances simplicity and efficiency.",
        ],
        "notes": "This project shows how theoretical structures become practical when they are matched carefully to the application's workload. The final search engine is a compact example of that principle.",
        "script": "To conclude, this project shows that the right data structures have a direct and measurable effect on performance. The final design is not only correct, but also optimized in ways that matter for real search workloads.",
    },
]


def set_slide_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title: str) -> None:
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.82))
    banner.fill.solid()
    banner.fill.fore_color.rgb = TITLE_COLOR
    banner.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.78), Inches(13.333), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_COLOR
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.14), Inches(12.2), Inches(0.5))
    text_frame = title_box.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = WHITE


def add_footer(slide, slide_number: int) -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(7.03), Inches(12.23), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(226, 232, 240)
    line.line.fill.background()

    footer = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(11.3), Inches(0.2))
    frame = footer.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = "UC Algo Phase 4 Presentation"
    p.font.name = "Aptos"
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED_COLOR

    num_box = slide.shapes.add_textbox(Inches(12.1), Inches(7.05), Inches(0.7), Inches(0.2))
    num_frame = num_box.text_frame
    num_frame.clear()
    p2 = num_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    p2.text = f"{slide_number:02d}"
    p2.font.name = "Aptos"
    p2.font.size = Pt(9)
    p2.font.color.rgb = MUTED_COLOR


def add_bullets(slide, bullets: list[str]) -> None:
    box = slide.shapes.add_textbox(Inches(0.9), Inches(1.45), Inches(11.8), Inches(4.8))
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.clear()

    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.bullet = True
        paragraph.space_after = Pt(12)
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(22)
        paragraph.font.color.rgb = BODY_COLOR


def add_code_box(slide, text: str, left: float, top: float, width: float, height: float, fill: RGBColor) -> None:
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = RGBColor(210, 218, 226)

    text_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.15), Inches(width - 0.4), Inches(height - 0.3))
    frame = text_box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = "Cascadia Mono"
    p.font.size = Pt(12.5)
    p.font.color.rgb = BODY_COLOR


def add_notes(slide, notes: str, script: str) -> None:
    notes_text = f"Speaker Notes:\n{notes}\n\nScript:\n{script}"
    slide.notes_slide.shapes[1].text = notes_text


def build_title_slide(prs: Presentation, slide_data: dict, slide_number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, TITLE_COLOR)

    accent_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.7), Inches(0.6), Inches(2.1), Inches(2.1))
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = ACCENT_COLOR
    accent_shape.fill.transparency = 0.18
    accent_shape.line.fill.background()

    accent_shape_2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.9), Inches(4.7), Inches(1.35), Inches(1.35))
    accent_shape_2.fill.solid()
    accent_shape_2.fill.fore_color.rgb = RGBColor(255, 255, 255)
    accent_shape_2.fill.transparency = 0.88
    accent_shape_2.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(9.4), Inches(1.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_r = title_p.add_run()
    title_r.text = slide_data["title"]
    title_r.font.name = "Aptos Display"
    title_r.font.size = Pt(30)
    title_r.font.bold = True
    title_r.font.color.rgb = WHITE

    subtitle_box = slide.shapes.add_textbox(Inches(0.82), Inches(3.1), Inches(9.2), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.clear()
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_r = subtitle_p.add_run()
    subtitle_r.text = slide_data["subtitle"]
    subtitle_r.font.name = "Aptos"
    subtitle_r.font.size = Pt(18)
    subtitle_r.font.color.rgb = RGBColor(225, 232, 240)

    meta_box = slide.shapes.add_textbox(Inches(0.82), Inches(5.62), Inches(5.2), Inches(0.9))
    meta_frame = meta_box.text_frame
    meta_frame.clear()
    meta_p = meta_frame.paragraphs[0]
    meta_r = meta_p.add_run()
    meta_r.text = "Phase 4 final report and presentation deck"
    meta_r.font.name = "Aptos"
    meta_r.font.size = Pt(14)
    meta_r.font.color.rgb = RGBColor(255, 244, 214)

    add_notes(slide, slide_data["notes"], slide_data["script"])
    add_footer(slide, slide_number)


def build_content_slide(prs: Presentation, slide_data: dict, slide_number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, SOFT_BG)
    add_title_bar(slide, slide_data["title"])

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.1), Inches(0.08), Inches(5.55))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_COLOR
    accent.line.fill.background()

    add_bullets(slide, slide_data["bullets"])

    if slide_data["title"] == "Algorithm Overview and Sample Output":
        pseudocode = (
            "PSEUDOCODE\n"
            "1. tokenize(query)\n"
            "2. if mode = AND then intersect posting lists\n"
            "3. else union posting lists\n"
            "4. for each candidate document:\n"
            "5.     score document with TF-IDF or BM25\n"
            "6.     push (doc, score) into top-k heap\n"
            "7. return results sorted by score"
        )
        sample_results = (
            "SAMPLE RESULTS\n"
            "Query: search index\n"
            "Result: doc-002\n\n"
            "Query: search graph\n"
            "Results: doc-001, doc-002, doc-004\n\n"
            "Observation: only candidate documents are scored; the full corpus is never sorted."
        )
        add_code_box(slide, pseudocode, 0.95, 1.55, 5.55, 4.95, RGBColor(255, 255, 255))
        add_code_box(slide, sample_results, 6.75, 1.55, 5.7, 4.95, RGBColor(255, 250, 242))

        left_label = slide.shapes.add_textbox(Inches(1.1), Inches(1.27), Inches(2.4), Inches(0.2))
        left_frame = left_label.text_frame
        left_frame.clear()
        p1 = left_frame.paragraphs[0]
        p1.text = "Algorithm Steps"
        p1.font.name = "Aptos"
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = ACCENT_COLOR

        right_label = slide.shapes.add_textbox(Inches(6.9), Inches(1.27), Inches(2.8), Inches(0.2))
        right_frame = right_label.text_frame
        right_frame.clear()
        p2 = right_frame.paragraphs[0]
        p2.text = "Sample Output"
        p2.font.name = "Aptos"
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = ACCENT_COLOR

        bullet_note = slide.shapes.add_textbox(Inches(0.95), Inches(6.6), Inches(11.7), Inches(0.35))
        note_frame = bullet_note.text_frame
        note_frame.clear()
        n = note_frame.paragraphs[0]
        n.text = "This layout separates the algorithm from its output so the presentation reads like a formal technical summary."
        n.font.name = "Aptos"
        n.font.size = Pt(10.5)
        n.font.color.rgb = MUTED_COLOR

    if slide_data["title"] == "Test Execution Results":
        command_box_text = (
            "AUTOMATED TEST COMMAND\n"
            "python -m unittest discover -s tests -p \"test_*.py\" -v"
        )
        results_box_text = (
            "EXECUTION SUMMARY\n"
            "Ran 17 tests in ~0.05s\n"
            "Status: OK (0 failures)\n\n"
            "Test modules:\n"
            "- test_data_structures.py\n"
            "- test_search_engine.py\n"
            "- test_phase3_scaling.py"
        )
        add_code_box(slide, command_box_text, 0.95, 1.62, 5.7, 2.1, RGBColor(255, 255, 255))
        add_code_box(slide, results_box_text, 0.95, 3.95, 5.7, 2.55, RGBColor(255, 250, 242))

        right_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.62), Inches(5.45), Inches(4.88))
        right_panel.fill.solid()
        right_panel.fill.fore_color.rgb = RGBColor(239, 244, 248)
        right_panel.line.color.rgb = RGBColor(199, 209, 219)

        right_text = slide.shapes.add_textbox(Inches(7.15), Inches(1.86), Inches(4.95), Inches(4.3))
        rf = right_text.text_frame
        rf.clear()
        rf.word_wrap = True
        p0 = rf.paragraphs[0]
        p0.text = "Why this matters"
        p0.font.name = "Aptos"
        p0.font.size = Pt(15)
        p0.font.bold = True
        p0.font.color.rgb = TITLE_COLOR

        for line in [
            "Validation is reproducible, not anecdotal.",
            "Both functional and scaling behaviors were tested.",
            "Results support claims made in the performance section.",
        ]:
            p = rf.add_paragraph()
            p.text = line
            p.level = 0
            p.bullet = True
            p.space_after = Pt(10)
            p.font.name = "Aptos"
            p.font.size = Pt(13)
            p.font.color.rgb = BODY_COLOR

    notes_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(5.45), Inches(4.35), Inches(1.12))
    notes_box.fill.solid()
    notes_box.fill.fore_color.rgb = RGBColor(239, 244, 248)
    notes_box.line.color.rgb = RGBColor(199, 209, 219)

    notes_text = slide.shapes.add_textbox(Inches(8.42), Inches(5.62), Inches(4.0), Inches(0.86))
    notes_frame = notes_text.text_frame
    notes_frame.clear()
    p = notes_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Speaker notes are stored in the deck notes section."
    r.font.name = "Aptos"
    r.font.size = Pt(10.5)
    r.font.color.rgb = MUTED_COLOR

    add_notes(slide, slide_data["notes"], slide_data["script"])
    add_footer(slide, slide_number)


def build_presentation() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_number, slide_data in enumerate(SLIDES, start=1):
        if slide_data.get("is_title"):
            build_title_slide(prs, slide_data, slide_number)
        else:
            build_content_slide(prs, slide_data, slide_number)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)


if __name__ == "__main__":
    build_presentation()
    print(f"Saved PowerPoint deck to {OUTPUT_PATH}")